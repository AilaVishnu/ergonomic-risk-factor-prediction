"""
Build a copy of the Final deck with Phase 4-7 results appended.

Creates deck/Ergonomic_Risk_Factor_Prediction_Project_Plan_WITH_RESULTS.pptx
from the existing Final.pptx, appends a Results & Findings section before
the Thank You slide, and saves.
"""

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC  = ROOT / "deck" / "Ergonomic_Risk_Factor_Prediction_Project_Plan_Final.pptx"
DST  = ROOT / "deck" / "Ergonomic_Risk_Factor_Prediction_Project_Plan_WITH_RESULTS.pptx"
FIGS = ROOT / "outputs" / "figures"


def find_thank_you_index(prs):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "Thank You" in shape.text_frame.text:
                return i
    return None


def add_image_slide(prs, title, image_filename, caption=None,
                    img_width=Inches(11)):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    if s.shapes.title is not None:
        s.shapes.title.text = title

    img_path = FIGS / image_filename
    if not img_path.exists():
        raise FileNotFoundError(img_path)

    pic = s.shapes.add_picture(str(img_path), Inches(1), Inches(1.5),
                               width=img_width)
    # Centre horizontally and cap height so it does not overrun the caption
    pic.left = int((prs.slide_width - pic.width) / 2)
    max_h = Inches(5.0)
    if pic.height > max_h:
        ratio = max_h / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)

    if caption:
        tb = s.shapes.add_textbox(Inches(0.5), Inches(6.7),
                                  prs.slide_width - Inches(1), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(14)
        r.font.italic = True
    return s


def add_text_slide(prs, title, body_lines):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    if s.shapes.title is not None:
        s.shapes.title.text = title

    tb = s.shapes.add_textbox(Inches(0.6), Inches(1.4),
                              prs.slide_width - Inches(1.2),
                              prs.slide_height - Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(18) if not line.startswith(" ") else Pt(16)
    return s


def add_section_header(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[2])  # Section Header
    if s.shapes.title is not None:
        s.shapes.title.text = title
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            ph.text_frame.text = subtitle
            break
    return s


def main():
    # If the destination file is open in PowerPoint, the copy will fail with
    # WinError 32. PowerPoint leaves a ~$ lock file next to the open file -
    # check for it and bail with a useful message instead of a stack trace.
    lock = DST.parent / f"~${DST.name}"
    if lock.exists():
        print(f"ERROR: {DST.name} is currently open in PowerPoint")
        print(f"       Close PowerPoint and re-run this script.")
        print(f"       (lock file present: {lock.name})")
        raise SystemExit(1)

    shutil.copy2(SRC, DST)
    prs = Presentation(DST)
    thank_you_idx = find_thank_you_index(prs)
    print(f"original slides : {len(prs.slides)}, Thank You at index {thank_you_idx}")

    # --- Results & Findings section ---

    add_section_header(prs, "Results & Findings",
                       "Phase 4-7 outputs from the 182-rider sample")

    add_image_slide(prs, "Sample demographics (n = 182)",
                    "demographics.png")

    add_image_slide(prs, "MSD prevalence by body area",
                    "nordic_prevalence.png",
                    caption="84.6% of riders reported pain in at least one body area in the last 12 months")

    add_image_slide(prs, "Risk factor distribution (all 6 factors)",
                    "risk_factor_distribution.png")

    add_image_slide(prs, "Discomfort prevalence by demographic group",
                    "discomfort_by_demographic.png",
                    caption="Strong age gradient: 72% (under 25) up to 100% (45+)")

    add_image_slide(prs, "Discomfort prevalence at each risk level",
                    "risk_vs_discomfort.png",
                    caption="Force shows the cleanest dose-response of the 5 survey-derived factors")

    add_image_slide(prs, "Correlation heatmap",
                    "correlation_heatmap.png",
                    caption="Workload, age, and job duration are the top correlations with discomfort")

    add_text_slide(prs, "ML model performance (Stage 2)", [
        "Best model per risk factor (5-fold CV, after GridSearchCV + SMOTE):",
        "",
        "Force           HistGradientBoosting        62%   AUC 71%   (42 features)",
        "Repetition      Random Forest               62%   AUC 73%   (41 features)",
        "Duration        Random Forest               61%   AUC 76%   (42 features)",
        "Vibration       Extra Trees                 58%   AUC 72%   (41 features)",
        "Contact Stress  Random Forest               60%   AUC 74%   (42 features)",
        "Posture         HistGradientBoosting        97%   AUC 98%   (63 features)",
        "",
        "5 survey-derived factors land inside the 60-80% published survey-based band.",
        "Posture uses real RULA + QEC observation inputs and reaches sensor-based range.",
    ])

    add_image_slide(prs, "Confusion matrices (best model per factor)",
                    "confusion_matrices.png")

    add_image_slide(prs, "ROC curves (one-vs-rest)",
                    "roc_curves.png",
                    caption="AUC > 0.70 for 5 of 6 factors")

    add_image_slide(prs, "Top features per model",
                    "feature_importance.png",
                    caption="workload_score, fatigue_score, and age_ord recur across factors")

    add_text_slide(prs, "Comparison with published benchmarks", [
        "Survey-based MSD-prediction studies typically reach 60-80% accuracy",
        "(Annals of Occupational and Environmental Medicine 2024, review of 130 studies).",
        "",
        "Sensor-based studies (IMU / EMG / computer vision) reach 90-99% because",
        "the inputs are direct physical signals.",
        "",
        "Our 5 survey-derived factors all land inside the survey-based band.",
        "",
        "Posture uses real RULA + QEC observation inputs (not survey proxies),",
        "so it sits in the sensor-based range at 97% / AUC 98%.",
    ])

    add_text_slide(prs, "Limitations", [
        "1. n = 182 self-report sample. Cross-sectional, single-region (Tamil Nadu).",
        "",
        "2. Posture per-rider linkage is approximate. RULA + QEC observations shared",
        "   no rider identifier with the survey - a severity-rank merge was used.",
        "   The 97% accuracy is the upper bound the linked data permits; a",
        "   per-rider observation study would likely settle slightly lower.",
        "",
        "3. Repetition binning was corrected (qcut -> fixed cuts [2.5, 3.75]) so",
        "   the worst real combo no longer ties on the tercile edge.",
        "   Stage-2 accuracy dropped 74% -> 62% but is now methodologically honest.",
        "",
        "4. Vibration's 58% accuracy reflects a real signal ceiling once vehicle_rank,",
        "   work_hours_num and vibration_index are all excluded to prevent leakage.",
        "",
        "5. Self-report biases apply to discomfort, fatigue and workload measures.",
    ])

    add_text_slide(prs, "Recommendations for rider safety", [
        "Duration (49% of riders in High band - largest population burden):",
        "   shift-length caps and mandatory breaks every 2 hours",
        "",
        "Force (only factor significantly associated with discomfort, p = 0.035):",
        "   limit carried-package weight; train safer lifting technique",
        "",
        "Contact Stress (handheld carriers reported less discomfort):",
        "   encourage bike-storage box or backpack over handheld delivery",
        "",
        "Age and job duration (strongest individual-level discomfort predictors):",
        "   workload tapering for riders over 12 months tenure or over 35 years old",
    ])

    # Move "Thank You" to the very end so the results sit before it
    if thank_you_idx is not None:
        xml_slides = prs.slides._sldIdLst
        slides_xml = list(xml_slides)
        ty = slides_xml[thank_you_idx]
        xml_slides.remove(ty)
        xml_slides.append(ty)

    prs.save(DST)
    print(f"saved: {DST.relative_to(ROOT)}")
    print(f"final slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
