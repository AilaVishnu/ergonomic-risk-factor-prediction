"""
Build deck/Review_Presentation.pptx - a 10-slide review deck sized for
the 7-minute IIITDM-SIES mid-review presentation.

Slide plan (max 10 slides, ~40 seconds each):
  01. Title (project title + student + institute + mentor)
  02. Problem & Objectives
  03. Two-stage methodology (pipeline flowchart)
  04. Data sources
  05. Six risk factors + Stage-1 thresholds
  06. Stage-2 ML pipeline (algorithms + CV setup)
  07. Statistical findings (NMQ prevalence, logistic regression)
  08. Model performance (Stage-2 metrics table)
  09. Web application demo
  10. Conclusions, recommendations, future work

Run:
    python src/build_review_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT  = ROOT / "deck" / "Review_Presentation.pptx"
FIGS = ROOT / "outputs" / "figures"
SHOT = ROOT / "outputs" / "app_screenshots"


# ----- Design tokens -----
NAVY    = RGBColor(0x0F, 0x1F, 0x3A)   # slide title, dark headings
INK     = RGBColor(0x1F, 0x2A, 0x40)   # body text
GREY    = RGBColor(0x6B, 0x72, 0x82)   # muted labels
ACCENT  = RGBColor(0x0F, 0x76, 0x6E)   # teal accent bar
LIGHT   = RGBColor(0xF7, 0xF8, 0xFA)   # section-band background
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)


def _set_font(run, name="Calibri", size=18, bold=False, italic=False,
              color=INK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_bg(slide, prs, colour=WHITE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = colour
    bg.shadow.inherit = False
    return bg


def _add_accent_bar(slide, prs, x=Inches(0.7), y=Inches(0.42),
                    w=Inches(0.35), h=Inches(0.05), colour=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.shadow.inherit = False


def _add_title(slide, text, x=Inches(0.7), y=Inches(0.85),
               w=Inches(11.8), size=32):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    _set_font(r, size=size, bold=True, color=NAVY)
    return tb


def _add_kicker(slide, text, x=Inches(0.7), y=Inches(0.55),
                w=Inches(11.8), size=10):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text.upper()
    _set_font(r, size=size, bold=True, color=ACCENT)
    r.font.color.rgb = ACCENT
    return tb


def _add_body_text(slide, text, x, y, w, h, size=16, bold=False,
                   italic=False, colour=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        _set_font(r, size=size, bold=bold, italic=italic, color=colour)
    return tb


def _add_bullets(slide, items, x, y, w, h, size=15, colour=INK):
    """Add a bullet list.  items: list of (text, indent_level) or plain text."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        # Bullet character prefix
        bullet_char = "•" if level == 0 else "–"
        r_b = p.add_run()
        r_b.text = f"{bullet_char}   "
        _set_font(r_b, size=size, bold=True, color=ACCENT)
        r_t = p.add_run()
        r_t.text = text
        _set_font(r_t, size=size, color=colour)
    return tb


def _add_footer(slide, prs, page_num, total):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.05),
                                  Inches(11.8), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("A Predictive ML Framework for Ergonomic Risk in "
              "Last-Mile Quick-Commerce Delivery Operations   |   "
              f"Slide {page_num}/{total}")
    _set_font(r, size=10, color=GREY)


def _add_image(slide, path, x, y, w=None, h=None, caption=None):
    if not Path(path).exists():
        return
    kwargs = {}
    if w: kwargs["width"] = w
    if h: kwargs["height"] = h
    pic = slide.shapes.add_picture(str(path), x, y, **kwargs)
    if caption:
        cap = slide.shapes.add_textbox(
            x, y + pic.height + Inches(0.05), pic.width, Inches(0.3)
        )
        cp = cap.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        r = cp.add_run()
        r.text = caption
        _set_font(r, size=10, italic=True, color=GREY)
    return pic


def _add_table(slide, header, rows, x, y, w, h,
               header_bg=NAVY, header_text=WHITE, size=12,
               col_widths=None):
    n_cols = len(header)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table

    # Explicit column widths (in Inches units, list of length n_cols).
    if col_widths is not None:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    # Header
    for i, col in enumerate(header):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = col
        _set_font(r, size=size, bold=True, color=header_text)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 else LIGHT
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            _set_font(r, size=size, color=INK)
    return tbl


# =========================================================================
# Slide builders
# =========================================================================

def slide_01_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(s, prs)

    # Accent bar top-left
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.7),
                             Inches(0.9), Inches(0.12))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.shadow.inherit = False

    # Kicker
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.95),
                              Inches(11.8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "IIITDM-SIES INTERNSHIP  |  MID-REVIEW PRESENTATION"
    _set_font(r, size=12, bold=True, color=ACCENT)

    # Big title
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.7),
                              Inches(11.8), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        "A Predictive Machine Learning Framework",
        "for Ergonomic Risk in Last-Mile",
        "Quick-Commerce Delivery Operations",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        _set_font(r, size=36, bold=True, color=NAVY)

    # Divider line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(4.7),
                              Inches(3.0), Inches(0.02))
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.shadow.inherit = False

    # Details block
    tb = s.shapes.add_textbox(Inches(0.7), Inches(4.95),
                              Inches(11.8), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True

    for kicker, value in [
        ("Presented by",  "AILA VISHNU VARDHAN"),
        ("Institute",     "Vidya Jyothi Institute of Technology"),
        ("Mentor",        "Dr. Arunachalam Muthiah"),
        ("Programme",     "IIITDM-SIES Internship  |  "
                          "School of Interdisciplinary Design and "
                          "Innovation (SIDI), IIITDM Kancheepuram"),
    ]:
        p = tf.add_paragraph()
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"{kicker.upper():<16}  "
        _set_font(r1, size=11, bold=True, color=GREY)
        r2 = p.add_run()
        r2.text = value
        _set_font(r2, size=15, color=INK)

    # First paragraph is blank; nuke it
    tf.paragraphs[0].text = ""

    # Date footer
    tb = s.shapes.add_textbox(Inches(0.7), Inches(7.0),
                              Inches(11.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "July 2026"
    _set_font(r, size=11, italic=True, color=GREY)


def slide_02_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "01  Background")
    _add_title(s, "Problem and Objectives")

    _add_body_text(
        s,
        "Last-mile quick-commerce delivery riders (Blinkit, Zepto) face "
        "prolonged awkward postures, high delivery rates, heavy carrying "
        "loads and vehicle vibration, leading to musculoskeletal disorders "
        "(MSDs). No single tool combines the six standard ergonomic risk "
        "factors into an automated per-rider screening.",
        Inches(0.7), Inches(1.85), Inches(11.8), Inches(1.5),
        size=15,
    )

    # Objectives
    _add_body_text(s, "Objectives",
                   Inches(0.7), Inches(3.6), Inches(11.8), Inches(0.4),
                   size=17, bold=True, colour=NAVY)

    _add_bullets(s, [
        "Build a two-stage pipeline (deterministic labels + supervised ML) "
        "that scores each rider on Force, Repetition, Posture, Duration, "
        "Contact Stress, and Vibration.",
        "Train seven candidate classifiers per factor with SMOTE + "
        "GridSearchCV inside 5-fold stratified cross-validation.",
        "Identify and correct methodological issues discovered during "
        "modelling (label leakage, boundary degeneracy in binning).",
        "Deploy an interactive Streamlit web application that accepts the "
        "full 36-item questionnaire and returns six colour-coded risk levels.",
    ],
    Inches(0.7), Inches(4.1), Inches(11.8), Inches(2.7), size=14)

    _add_footer(s, prs, 2, 10)


def slide_03_methodology(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "02  Methodology")
    _add_title(s, "Two-Stage Pipeline")

    # Left column: text
    _add_bullets(s, [
        "Stage 1 (deterministic).  Standard ergonomic methods produce "
        "Low / Medium / High labels per factor.",
        ("Borg CR10 operational thresholds for Force", 1),
        ("RULA action levels collapsed into three bands for Posture", 1),
        ("Fixed cuts / sample terciles for the remaining four", 1),
        "Stage 2 (supervised ML).  A classifier per factor learns the "
        "Stage-1 label from the rider profile, with per-target exclusions "
        "to prevent trivial leakage.",
        "Two-stage split lets an ergonomist audit Stage 1 by hand, "
        "independently of the ML review.",
    ],
    Inches(0.7), Inches(1.85), Inches(6.6), Inches(5), size=13)

    # Right column: pipeline flowchart figure
    flowchart = FIGS / "methodology_flowchart.png"
    if flowchart.exists():
        _add_image(s, flowchart, Inches(7.6), Inches(2.3),
                   w=Inches(5.5),
                   caption="Pipeline flowchart")

    _add_footer(s, prs, 3, 10)


def slide_04_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "03  Data")
    _add_title(s, "Data Sources")

    # Two side-by-side cards
    def _card(x, y, w, h, kicker_text, title, body_lines):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = RGBColor(0xE0, 0xE3, 0xE9)
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        _add_body_text(s, kicker_text.upper(),
                       x + Inches(0.3), y + Inches(0.2),
                       w - Inches(0.6), Inches(0.3),
                       size=10, bold=True, colour=ACCENT)
        _add_body_text(s, title,
                       x + Inches(0.3), y + Inches(0.55),
                       w - Inches(0.6), Inches(0.5),
                       size=17, bold=True, colour=NAVY)
        _add_bullets(s, body_lines,
                     x + Inches(0.3), y + Inches(1.15),
                     w - Inches(0.6), h - Inches(1.4),
                     size=12)

    _card(Inches(0.7), Inches(1.9), Inches(6.15), Inches(4.8),
          "n = 182",
          "Rider Survey (CSV)",
          ["36 questionnaire items across 5 modules.",
           ("Q1-17: demographics, work pattern, lifestyle", 1),
           ("Q18: NMQ 12-month pain (9 body areas)", 1),
           ("Q19: 7-day discomfort (4 areas)", 1),
           ("Q20-24: discomfort follow-ups", 1),
           ("Q25-30: NASA-TLX 0-100 sliders", 1),
           ("Q31-36: Borg CR10 0-10 sliders", 1),
           "Self-administered; single time-point.",
          ])

    _card(Inches(7.15), Inches(1.9), Inches(5.65), Inches(4.8),
          "n = 182",
          "Posture Observations (xlsx)",
          ["Ergonomic observation records.",
           ("11 RULA components", 1),
           ("3 RULA Table A/B/C scores", 1),
           ("8 QEC region + exposure scores", 1),
           "Merged into the survey via severity-rank pairing.",
           ("Ranks riders by NMQ + fatigue + hours", 1),
           ("Ranks observations by RULA Table C", 1),
           ("Matches rank-to-rank one-to-one", 1),
          ])

    _add_footer(s, prs, 4, 10)


def slide_05_factors(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "04  Stage 1")
    _add_title(s, "Six Risk Factors & Stage-1 Distribution")

    header = ["Factor", "Method", "Threshold rule", "Low", "Med", "High"]
    rows = [
        ["Force",          "Borg CR10 (lifting)",     "0-3 / 4-6 / 7-10 dph",       "90",  "57",  "35"],
        ["Repetition",     "Deliveries per hour",     "cut at 2.5 & 3.75 dph",      "26",  "82",  "74"],
        ["Duration",       "Continuous hours",        "<= 5 / 6-7 / > 7 hrs",       "37",  "56",  "89"],
        ["Vibration",      "vehicle_rank * hours",    "sample tercile",              "67",  "68",  "47"],
        ["Contact Stress", "carry x hours x age",     "sample tercile",              "68",  "58",  "56"],
        ["Posture",        "RULA Table C",            "1-2 / 3-4 / 5+ (AL1-4)",     "0",   "29",  "153"],
    ]
    _add_table(s, header, rows,
               Inches(0.7), Inches(1.9), Inches(11.8), Inches(3.5), size=12)

    _add_body_text(
        s,
        "Posture, Duration, and Repetition are the three factors where the "
        "High band dominates.  Posture has no Low observations because "
        "the training-data minimum RULA Table C score is 3 (Medium).",
        Inches(0.7), Inches(5.7), Inches(11.8), Inches(1),
        size=13, colour=GREY,
    )

    _add_footer(s, prs, 5, 10)


def slide_06_ml(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "05  Stage 2")
    _add_title(s, "Machine-Learning Pipeline")

    # Left column
    _add_body_text(s, "Seven candidate classifiers",
                   Inches(0.7), Inches(1.85), Inches(6.3), Inches(0.4),
                   size=15, bold=True, colour=NAVY)
    _add_bullets(s, [
        "Logistic Regression (L2, balanced weights)",
        "Decision Tree, Random Forest, Extra Trees",
        "Hist Gradient Boosting, XGBoost",
        "Stacking Classifier (RF + ET + XGB + HGB)",
    ],
    Inches(0.7), Inches(2.3), Inches(6.3), Inches(2.4), size=13)

    _add_body_text(s, "Cross-validation and tuning",
                   Inches(0.7), Inches(4.5), Inches(6.3), Inches(0.4),
                   size=15, bold=True, colour=NAVY)
    _add_bullets(s, [
        "StratifiedKFold(5, shuffle=True, random_state=42)",
        "SMOTE inside every training fold (imblearn.Pipeline)",
        "GridSearchCV on macro F1 for hyperparameters",
        "Best model per target refit on full sample and saved as .pkl",
    ],
    Inches(0.7), Inches(4.95), Inches(6.3), Inches(2), size=13)

    # Right column
    _add_body_text(s, "Per-target feature exclusions",
                   Inches(7.4), Inches(1.85), Inches(5.4), Inches(0.4),
                   size=15, bold=True, colour=NAVY)

    header = ["Target", "Excluded features", "Feats"]
    rows = [
        ["Force",          "force_exertion, force_x_age",              "42"],
        ["Repetition",     "deliveries_num, hours, deliv_x_days",       "41"],
        ["Duration",       "work_hours_num, vibration_index",           "42"],
        ["Vibration",      "vibration_index, vehicle, hours",           "41"],
        ["Contact Stress", "carry, work_hours_num",                     "42"],
        ["Posture",        "(RULA Tables A/B/C only)",                  "63"],
    ]
    _add_table(s, header, rows,
               Inches(7.4), Inches(2.3), Inches(5.4), Inches(3), size=10)

    _add_body_text(
        s,
        "Prevents label leakage: each model has to learn from the rest of "
        "the rider profile instead of memorising the Stage-1 rule.",
        Inches(7.4), Inches(5.5), Inches(5.4), Inches(1),
        size=11, italic=True, colour=GREY,
    )

    _add_footer(s, prs, 6, 10)


def slide_07_stats(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "06  Findings")
    _add_title(s, "Statistical Findings")

    # ------------------------------------------------------------------
    # 2 x 2 quadrant layout
    #   TL: NMQ prevalence figure          TR: Logistic regression predictors
    #   BL: Chi-square all 6 factors       BR: Mann-Whitney 7-day discomfort
    # ------------------------------------------------------------------

    # --- Top-Left: NMQ figure (aspect 1.818).  Height-limited to stay
    #     inside the top-left quadrant (row 1: y 1.9 -> 4.35).
    nmq = FIGS / "nordic_prevalence.png"
    if nmq.exists():
        img_h = 2.1
        img_w = img_h * 1.818                          # 3.82"
        img_x = 0.7 + (5.9 - img_w) / 2                # centered in quad
        _add_image(s, nmq, Inches(img_x), Inches(1.95),
                   h=Inches(img_h),
                   caption="NMQ 12-month pain prevalence (n = 182)")

    # --- Top-Right: Logistic regression -------------------------------
    _add_body_text(s, "Logistic regression  |  significant predictors of any 12-month pain",
                   Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.35),
                   size=11, bold=True, colour=NAVY)

    lr_header = ["Predictor",           "OR",   "95% CI",         "p"]
    lr_rows = [
        ["Age (per band)",              "3.58", "1.70 - 7.56",    "0.0008"],
        ["Job duration (per band)",     "2.89", "1.54 - 5.41",    "0.001"],
        ["Fatigue score",               "1.43", "1.13 - 1.80",    "0.003"],
        ["Income (per band)",           "2.00", "1.25 - 3.20",    "0.004"],
        ["Workload score",              "1.06", "1.02 - 1.09",    "0.0005"],
        ["Education (per band)",        "0.33", "0.12 - 0.94",    "0.04"],
    ]
    _add_table(s, lr_header, lr_rows,
               Inches(7.0), Inches(2.35), Inches(5.8), Inches(2.0),
               size=10)

    # --- Bottom-Left: Chi-square all 6 risk factors -------------------
    _add_body_text(s, "Chi-square  |  Stage-1 risk factor vs 12-month pain",
                   Inches(0.7), Inches(4.55), Inches(5.9), Inches(0.35),
                   size=11, bold=True, colour=NAVY)

    cs_header = ["Risk factor",     "chi2",   "df", "p",       "Sig."]
    cs_rows = [
        ["Posture",         "45.67", "1", "< 0.001",  "yes"],
        ["Repetition",      "8.62",  "2", "0.014",    "yes"],
        ["Force",           "6.72",  "2", "0.035",    "yes"],
        ["Duration",        "0.62",  "2", "0.733",    "-"],
        ["Vibration",       "0.60",  "2", "0.741",    "-"],
        ["Contact stress",  "0.54",  "2", "0.762",    "-"],
    ]
    _add_table(s, cs_header, cs_rows,
               Inches(0.7), Inches(5.0), Inches(5.9), Inches(1.65),
               size=10,
               col_widths=[Inches(2.0), Inches(0.9),
                           Inches(0.6), Inches(1.4), Inches(1.0)])

    # --- Bottom-Right: Mann-Whitney 7-day discomfort ------------------
    _add_body_text(s, "Mann-Whitney U  |  predictors of 7-day discomfort",
                   Inches(7.0), Inches(4.55), Inches(5.8), Inches(0.35),
                   size=11, bold=True, colour=NAVY)

    mw_header = ["Feature",       "no pain", "pain",  "U",     "p"]
    mw_rows = [
        ["Age (band)",             "0",   "1",   "1228",  "0.0001"],
        ["Tenure (band)",          "0",   "1",   "1262",  "0.0002"],
        ["Workload score",         "42",  "50",  "1291",  "0.0007"],
        ["Income (band)",          "1",   "2",   "1411",  "0.0023"],
        ["Fatigue score",          "3.8", "4.6", "1474",  "0.0078"],
        ["Force (Borg)",           "3",   "4",   "1585",  "0.025"],
    ]
    _add_table(s, mw_header, mw_rows,
               Inches(7.0), Inches(5.0), Inches(5.8), Inches(1.65),
               size=10,
               col_widths=[Inches(2.1), Inches(0.9),
                           Inches(0.8), Inches(1.0), Inches(1.0)])

    # --- Bottom summary line ------------------------------------------
    _add_body_text(
        s,
        "Posture, Repetition, and Force reach chi-square significance;  "
        "age, tenure, income, workload, and fatigue drive it at the "
        "individual level.",
        Inches(0.7), Inches(6.75), Inches(12.1), Inches(0.25),
        size=10, italic=True, colour=GREY,
    )

    _add_footer(s, prs, 7, 10)


def slide_08_performance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "07  Model Performance")
    _add_title(s, "Best Model per Factor")

    header = ["Factor",         "Best model",              "Acc.", "F1 macro", "AUC macro", "Features"]
    rows = [
        ["Force",           "HistGradientBoosting",  "62%",  "57%",  "71%",  "42"],
        ["Repetition",      "Random Forest",         "62%",  "57%",  "73%",  "41"],
        ["Duration",        "Random Forest",         "61%",  "58%",  "76%",  "42"],
        ["Vibration",       "Extra Trees",           "58%",  "57%",  "72%",  "41"],
        ["Contact Stress",  "Random Forest",         "60%",  "59%",  "74%",  "42"],
        ["Posture",         "HistGradientBoosting",  "97%",  "95%",  "98%",  "63"],
    ]
    _add_table(s, header, rows,
               Inches(0.7), Inches(1.9), Inches(11.8), Inches(3.4), size=13)

    _add_body_text(
        s,
        "The five survey-derived factors land at 58-62% accuracy with macro "
        "AUC 71-76%.  Posture reaches 97% / 98% because it is the only "
        "model that receives real observation inputs (11 RULA components "
        "+ 8 QEC scores) on top of the survey features.",
        Inches(0.7), Inches(5.5), Inches(11.8), Inches(1.4),
        size=13,
    )

    _add_footer(s, prs, 8, 10)


def slide_09_webapp(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "08  Deliverable")
    _add_title(s, "Interactive Web Application")

    # 3 x 2 grid of six app screenshots.  Aspect 1.455 -> tile 3.2" wide
    # gives 2.20" tall.  Row 1 y=1.82, row 2 y=4.55, footer at 7.05.
    tile_w_in = 3.2
    tile_w    = Inches(tile_w_in)
    tile_h_v  = tile_w_in / 1.455                        # 2.199"

    n_cols  = 3
    total_w = 3 * tile_w_in + 2 * 0.15                   # 9.9"
    left    = (13.333 - total_w) / 2                     # 1.72"

    row_y = [Inches(1.82), Inches(4.55)]

    shots = [
        ("web_01_home.png",              "Home"),
        ("web_02_assessment_top.png",    "Assessment  |  demographics"),
        ("web_03_assessment_nmq.png",    "Assessment  |  NMQ 12-month pain"),
        ("web_06_results.png",           "Results  |  per-factor risk"),
        ("web_07_methodology.png",       "Methodology"),
        ("web_08_about.png",             "About"),
    ]

    for i, (fname, cap) in enumerate(shots):
        r, c = divmod(i, n_cols)
        x = Inches(left + c * (tile_w_in + 0.15))
        y = row_y[r]
        path = SHOT / fname
        if not path.exists():
            continue
        pic = s.shapes.add_picture(str(path), x, y, width=tile_w)
        # Thin border around each tile so screenshots read as tiles
        pic.line.color.rgb = RGBColor(0xE0, 0xE3, 0xE9)
        pic.line.width = Pt(0.5)
        # Caption directly under the tile
        cap_tb = s.shapes.add_textbox(
            x, y + Inches(tile_h_v + 0.03),
            tile_w, Inches(0.22),
        )
        p = cap_tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r_run = p.add_run()
        r_run.text = cap
        _set_font(r_run, size=10, italic=True, color=GREY)

    _add_footer(s, prs, 9, 10)


def slide_10_conclusions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(s, prs)
    _add_accent_bar(s, prs)
    _add_kicker(s, "09  Conclusions")
    _add_title(s, "Conclusions & Future Work")

    # Left column: key contributions
    _add_body_text(s, "Key contributions",
                   Inches(0.7), Inches(1.85), Inches(6.3), Inches(0.4),
                   size=15, bold=True, colour=NAVY)
    _add_bullets(s, [
        "Two-stage pipeline auditable by ergonomists and ML reviewers "
        "independently.",
        "Six per-factor classifiers with honest 5-fold CV metrics.",
        "Two methodological corrections openly disclosed "
        "(Repetition qcut boundary, Duration leakage via vibration_index).",
        "Multi-page Streamlit web app for real-time screening.",
    ],
    Inches(0.7), Inches(2.3), Inches(6.3), Inches(3.2), size=12)

    # Right column: recommendations
    _add_body_text(s, "Platform-level recommendations",
                   Inches(7.1), Inches(1.85), Inches(5.7), Inches(0.4),
                   size=15, bold=True, colour=NAVY)
    _add_bullets(s, [
        "Cap daily hours (49% of sample > 8 h/day).",
        "Posture training and equipment review "
        "(84% at RULA action level 5+).",
        "Prefer bike-storage carriers over handheld bags.",
        "Age-targeted MSD screening for 36+ cohort.",
        "Platform-level workload management.",
    ],
    Inches(7.1), Inches(2.3), Inches(5.7), Inches(3.2), size=12)

    # Future work strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.7), Inches(5.8),
                               Inches(12.1), Inches(1.05))
    strip.fill.solid()
    strip.fill.fore_color.rgb = LIGHT
    strip.line.color.rgb = RGBColor(0xE0, 0xE3, 0xE9)
    strip.line.width = Pt(0.75)
    strip.shadow.inherit = False

    _add_body_text(s, "Future work",
                   Inches(0.9), Inches(5.9), Inches(11.8), Inches(0.4),
                   size=11, bold=True, colour=ACCENT)
    _add_body_text(
        s,
        "Longitudinal follow-up study  |  per-rider RULA observations "
        "(eliminate the severity-rank merge)  |  wearable-accelerometer "
        "Vibration proxy  |  periodic re-training as the workforce evolves.",
        Inches(0.9), Inches(6.15), Inches(12), Inches(0.7),
        size=12, colour=INK,
    )

    _add_footer(s, prs, 10, 10)


# =========================================================================
# Build
# =========================================================================

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_methodology(prs)
    slide_04_data(prs)
    slide_05_factors(prs)
    slide_06_ml(prs)
    slide_07_stats(prs)
    slide_08_performance(prs)
    slide_09_webapp(prs)
    slide_10_conclusions(prs)

    # Metadata
    cp = prs.core_properties
    cp.author = "AILA VISHNU VARDHAN"
    cp.title = ("Ergonomic Risk in Last-Mile Quick-Commerce Delivery — "
                "IIITDM-SIES Review")
    cp.comments = ""

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    p = build()
    print(f"saved {p.relative_to(ROOT)}")
    print(f"slides: 10")
