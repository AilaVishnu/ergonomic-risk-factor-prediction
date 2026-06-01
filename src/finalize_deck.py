"""
Phase 9 — Finalize the deck.

Applies the inconsistency fixes documented in `docs/development_plan.md`
under the heading "Inconsistencies still in the deck":

- **Slides 10 and 32** — rename
    "Predict Awkward Posture Risk" -> "Predict Posture Risk"
    "Predict Static Posture Risk"  -> "Predict Duration Risk"
  to match slide 3 and the supervisor's decision.

- **Slide 37** (the detailed Inputs/Outputs slide):
    - rename outputs "Static Posture Risk" -> "Duration Risk"
                     "Awkward Posture Risk" -> "Posture Risk"
    - "posture scores = RULA" -> "posture scores = RULA & QEC"
      (to match the method declared on slide 4)
    - fix the stray closing quote in `Working Hours per Day"`
    - DROP the line `- discomfort indicators = Nordic pain columns (...)`
      from the input list — discomfort is the prediction *target*, not a
      predictor (this is the corrected-plan rule against circularity).

- **Slide 36** — deleted entirely (it is the less-detailed duplicate of
  slide 37). Slide 37 becomes slide 36 after the delete.

A timestamped backup of the .pptx is written into `deck/archive/`
before any change is applied.
"""

from datetime import datetime
from pathlib import Path
import shutil

from pptx import Presentation

DECK = (
    Path(__file__).resolve().parents[1]
    / "deck"
    / "Ergonomic_Risk_Factor_Prediction_Project_Plan_Final.pptx"
)
ARCHIVE = DECK.parent / "archive"


def _replace_in_slide(slide, old, new):
    """Replace `old` with `new` in any text frame on the slide.

    Joins all runs in a paragraph, replaces, then writes the result back
    to the first run while clearing the rest. This preserves the
    paragraph-level formatting (font, size, color) of the first run."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(run.text for run in para.runs)
            if old in full:
                if para.runs:
                    para.runs[0].text = full.replace(old, new)
                    for run in para.runs[1:]:
                        run.text = ""
                    n += 1
    return n


def _delete_paragraph_containing(slide, substring):
    """Delete every paragraph that contains `substring`."""
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in list(shape.text_frame.paragraphs):
            full = "".join(run.text for run in para.runs)
            if substring in full:
                p = para._p
                p.getparent().remove(p)
                n += 1
    return n


def main():
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    ARCHIVE.mkdir(exist_ok=True)
    backup = (
        ARCHIVE
        / f"Ergonomic_Risk_Factor_Prediction_Project_Plan_BEFORE_phase9_{ts}.pptx"
    )
    shutil.copy2(DECK, backup)
    print(f"Backup written: deck/archive/{backup.name}")

    prs = Presentation(str(DECK))
    print(f"Loaded {len(prs.slides)} slides")

    # Slides 10 and 32 — rename Awkward / Static
    for slide_idx in (9, 31):
        slide = prs.slides[slide_idx]
        a = _replace_in_slide(
            slide, "Predict Awkward Posture Risk", "Predict Posture Risk"
        )
        b = _replace_in_slide(
            slide, "Predict Static Posture Risk", "Predict Duration Risk"
        )
        print(
            f"  Slide {slide_idx + 1}: renamed Awkward->Posture ({a}), "
            f"Static->Duration ({b})"
        )

    # Slide 37 — rename outputs, &QEC, fix quote, drop discomfort-input
    s37 = prs.slides[36]
    print("  Slide 37 edits:")
    print(
        "    Awkward Posture Risk -> Posture Risk :",
        _replace_in_slide(s37, "Awkward Posture Risk", "Posture Risk"),
    )
    print(
        "    Static Posture Risk  -> Duration Risk :",
        _replace_in_slide(s37, "Static Posture Risk", "Duration Risk"),
    )
    print(
        "    posture scores = RULA -> RULA & QEC   :",
        _replace_in_slide(
            s37, "posture scores = RULA", "posture scores = RULA & QEC"
        ),
    )
    print(
        "    fix stray closing quote on Working Hours :",
        _replace_in_slide(
            s37, 'Working Hours per Day"', "Working Hours per Day"
        ),
    )
    print(
        "    drop discomfort-indicators input line   :",
        _delete_paragraph_containing(
            s37, "discomfort indicators = Nordic pain"
        ),
    )

    # Slide 36 — delete (duplicate of 37)
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    xml_slides.remove(slide_ids[35])
    print("  Slide 36 deleted (duplicate of slide 37)")

    prs.save(str(DECK))
    print(
        f"Saved: {DECK.name}  ({len(prs.slides)} slides total after delete)"
    )


if __name__ == "__main__":
    main()
